# Copyright (c) 2025 ETH Zurich.
#                    All rights reserved.
#
# Use of this source code is governed by a BSD-style license that can be
# found in the LICENSE file.
#
# Main authors: Lorenzo Paleari
#               Jón Gunnar Hannesson
#
# Most of the code below is from the Microsoft Autogen repository.
# https://github.com/microsoft/autogen/blob/gaia_multiagent_v01_march_1st/autogen/mdconvert.py
#
# Copyright (c) Microsoft Corporation.

import copy
import html
import json
import logging
import mimetypes
import os
import re
import tempfile
import traceback
import xml.etree.ElementTree as ET
from abc import ABC, abstractmethod
from time import time
from typing import List, Optional, Union
from urllib.parse import parse_qs, urlparse

import mammoth
import markdownify
import pandas as pd
import pdfminer
import pdfminer.high_level
import pptx
import puremagic
import requests
from bs4 import BeautifulSoup
from langchain_community.callbacks import get_openai_callback
from openai import OpenAI
from youtube_transcript_api import YouTubeTranscriptApi

from kgot.utils import UsageStatistics

logger = logging.getLogger("Controller.MdConverter")

class DocumentConverterResult:
    """The result of converting a document to text."""

    def __init__(self, title: Union[str, None] = None, text_content: str = ""):
        self.title = title
        self.text_content = text_content


class DocumentConverter(ABC):
    @abstractmethod
    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        raise NotImplementedError()

class AudioConverter(DocumentConverter):
    usage_statistics = None
    def __init__(self, usage_statistics: Optional[UsageStatistics] = None):
        super().__init__()
        if usage_statistics is not None:
            self.usage_statistics = usage_statistics

    # Our audio transcriber, only uses whisper-1
    def transcribe_audio_our(self, audio_path) -> str:
        openai_obj = OpenAI(
            api_key=os.environ.get("OPENAI_API_KEY"),
            organization=os.environ.get("OPENAI_API_ORGANIZATION")
        )
        model = "whisper-1"
        with open(audio_path, "rb") as audio_file:
            if self.usage_statistics is not None:
                with get_openai_callback() as cb:
                    time_before = time()
                    transcription = openai_obj.audio.transcriptions.create(
                        model=model,
                        file=audio_file
                    )
                    time_after = time()
                    self.usage_statistics.log_statistic("AudioTranscriptionLoader.transcribe_audio",
                                                    time_before, time_after,
                                                    model,
                                                    cb.prompt_tokens, cb.completion_tokens, round(cb.total_cost, 6))
            else:
                transcription = openai_obj.audio.transcriptions.create(
                        model=model,
                        file=audio_file
                    )
        return transcription.text
    
    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not an audio file
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".wav", ".mp3", ".flac", ".m4a"]:
            return None
        result = self.transcribe_audio_our(local_path)
        return DocumentConverterResult(
            title=None,
            text_content=result,
        )

# All functions below are copied from the Beating the GAIA benchmark with Transformers Agents repository.
# https://github.com/aymeric-roucher/GAIA/blob/main/scripts/tools/mdconvert.py
#
# Copyright 2024 Aymeric Roucher
#
# Some minor modification to adapt the code to our needs were performed.
class PlainTextConverter(DocumentConverter):
    """Anything with content type text/plain"""

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension == "":
            return None

        encoding = kwargs.get("encoding", "utf-8")

        text_content = ""
        with open(local_path, "rt", encoding=encoding) as fh:
            text_content = fh.read()

        return DocumentConverterResult(
            title=None,
            text_content=text_content,
        )


class HtmlConverter(DocumentConverter):
    """Anything with content type text/html"""

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not html
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None
        encoding = kwargs.get("encoding", "utf-8")
        result = None
        with open(local_path, "rt", encoding=encoding) as fh:
            result = self._convert(fh.read())

        return result

    def _convert(self, html_content) -> Union[None, DocumentConverterResult]:
        """Helper function that converts and HTML string."""

        # Parse the string
        soup = BeautifulSoup(html_content, "html.parser")

        # Remove javascript and style blocks
        for script in soup(["script", "style"]):
            script.extract()

        # Print only the main content
        body_elm = soup.find("body")
        webpage_text = ""
        if body_elm:
            webpage_text = markdownify.MarkdownConverter().convert_soup(body_elm)
        else:
            webpage_text = markdownify.MarkdownConverter().convert_soup(soup)

        return DocumentConverterResult(
            title=None if soup.title is None else soup.title.string,
            text_content=webpage_text,
        )


class YouTubeConverter(DocumentConverter):
    """Handle YouTube specially, focusing on the video title, description, and transcript."""

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not YouTube
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None
        url = kwargs.get("url", "")
        if not url.startswith("https://www.youtube.com/watch?"):
            return None
        encoding = kwargs.get("encoding", "utf-8")
        # Parse the file
        soup = None
        with open(local_path, "rt", encoding=encoding) as fh:
            soup = BeautifulSoup(fh.read(), "html.parser")

        # Read the meta tags
        metadata = {"title": soup.title.string}
        for meta in soup(["meta"]):
            for a in meta.attrs:
                if a in ["itemprop", "property", "name"]:
                    metadata[meta[a]] = meta.get("content", "")
                    break

        # We can also try to read the full description. This is more prone to breaking, since it reaches into the page implementation
        try:
            for script in soup(["script"]):
                content = script.text
                if "ytInitialData" in content:
                    lines = re.split(r"\r?\n", content)
                    obj_start = lines[0].find("{")
                    obj_end = lines[0].rfind("}")
                    if obj_start >= 0 and obj_end >= 0:
                        data = json.loads(lines[0][obj_start : obj_end + 1])
                        attrdesc = self._findKey(data, "attributedDescriptionBodyText")
                        if attrdesc:
                            metadata["description"] = attrdesc["content"]
                    break
        except:
            pass

        # Start preparing the page
        webpage_text = "# YouTube\n"

        title = self._get(metadata, ["title", "og:title", "name"])
        if title:
            webpage_text += f"\n## {title}\n"

        stats = ""
        views = self._get(metadata, ["interactionCount"])
        if views:
            stats += f"- **Views:** {views}\n"

        keywords = self._get(metadata, ["keywords"])
        if keywords:
            stats += f"- **Keywords:** {keywords}\n"

        runtime = self._get(metadata, ["duration"])
        if runtime:
            stats += f"- **Runtime:** {runtime}\n"

        if len(stats) > 0:
            webpage_text += f"\n### Video Metadata\n{stats}\n"

        description = self._get(metadata, ["description", "og:description"])
        if description:
            webpage_text += f"\n### Description\n{description}\n"

        transcript_text = ""
        parsed_url = urlparse(url)
        params = parse_qs(parsed_url.query)

        video_id = params["v"][0]
        # Must be a single transcript.
        print("VIDDDD ID:", video_id)
        transcript = YouTubeTranscriptApi.get_transcript(video_id)
        transcript_text = " ".join([part["text"] for part in transcript])
        # Alternative formatting:
        # formatter = TextFormatter()
        # formatter.format_transcript(transcript)
        if transcript_text:
            webpage_text += f"\n### Transcript\n{transcript_text}\n"

        return DocumentConverterResult(
            title=title if title else soup.title.string,
            text_content=webpage_text,
        )

    def _get(self, json, keys, default=None):
        for k in keys:
            if k in json:
                return json[k]
        return default

    def _findKey(self, json, key):
        if isinstance(json, list):
            for elm in json:
                ret = self._findKey(elm, key)
                if ret is not None:
                    return ret
        elif isinstance(json, dict):
            for k in json:
                if k == key:
                    return json[k]
                else:
                    ret = self._findKey(json[k], key)
                    if ret is not None:
                        return ret
        return None


class PdfConverter(DocumentConverter):
    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a PDF
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pdf":
            return None

        return DocumentConverterResult(
            title=None,
            text_content=pdfminer.high_level.extract_text(local_path),
        )


class DocxConverter(HtmlConverter):
    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a DOCX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".docx":
            return None

        result = None
        with open(local_path, "rb") as docx_file:
            result = mammoth.convert_to_html(docx_file)
            html_content = result.value
            result = self._convert(html_content)

        return result


class XlsxConverter(HtmlConverter):
    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a XLSX
        extension = kwargs.get("file_extension", "")

        if extension.lower() not in [".xlsx", ".xls"]:
            return None

        sheets = pd.read_excel(local_path, sheet_name=None)
        md_content = ""
        for s in sheets:
            md_content += f"## {s}\n"
            html_content = sheets[s].to_html(index=False)
            md_content += self._convert(html_content).text_content.strip() + "\n\n"

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )


class XmlConverter(DocumentConverter):
    def convert(self, local_path, **kwargs) -> None | DocumentConverterResult:
        # Parse the XML string
        extension = kwargs.get("file_extension", "")

        if extension.lower() not in [".xml"]:
            return None
        encoding = kwargs.get("encoding", "utf-8")

        xml_string = ""
        with open(local_path, "rt", encoding=encoding) as fh:
            xml_string = fh.read()
        
        def extract_table_from_html_like(xml_root):
            table = xml_root.find('.//table')
            if table is None:
                raise ValueError("No table found in the XML")

            headers = [th.text for th in table.find('thead').findall('th')]
            rows = [[td.text for td in tr.findall('td')] for tr in table.find('tbody').findall('tr')]
            
            # Create markdown table
            markdown = '| ' + ' | '.join(headers) + ' |\n'
            markdown += '| ' + ' | '.join(['---'] * len(headers)) + ' |\n'
            for row in rows:
                markdown += '| ' + ' | '.join(row) + ' |\n'

        def extract_table_from_wordml(xml_root, namespaces):
            # Parse the XML content
            root = xml_root
            namespace = {'w': 'http://schemas.microsoft.com/office/word/2003/wordml'}
            
            # Extract text content
            body = root.find('w:body', namespace)
            paragraphs = body.findall('.//w:p', namespace)
            text_content = []
            for para in paragraphs:
                texts = para.findall('.//w:t', namespace)
                for text in texts:
                    text_content.append(text.text)
            
            return '\n'.join(text_content)

        # Parse the XML string
        root = ET.fromstring(xml_string)
        namespaces = {'w': 'http://schemas.microsoft.com/office/word/2003/wordml'}
        
        if root.tag.endswith('wordDocument'):
            markdown = extract_table_from_wordml(root, namespaces)
        else:
            markdown = extract_table_from_html_like(root)

        return DocumentConverterResult(
            title=None,
            text_content=markdown.strip(),
        )


class PptxConverter(HtmlConverter):
    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a PPTX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pptx":
            return None

        md_content = ""

        presentation = pptx.Presentation(local_path)
        slide_num = 0
        for slide in presentation.slides:
            slide_num += 1

            md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"

            title = slide.shapes.title
            for shape in slide.shapes:
                # Pictures
                if self._is_picture(shape):
                    # https://github.com/scanny/python-pptx/pull/512#issuecomment-1713100069
                    alt_text = ""
                    try:
                        alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                    except:
                        pass

                    # A placeholder name
                    filename = re.sub(r"\W", "", shape.name) + ".jpg"
                    # try:
                    #    filename = shape.image.filename
                    # except:
                    #    pass

                    md_content += "\n![" + (alt_text if alt_text else shape.name) + "](" + filename + ")\n"

                # Tables
                if self._is_table(shape):
                    html_table = "<html><body><table>"
                    first_row = True
                    for row in shape.table.rows:
                        html_table += "<tr>"
                        for cell in row.cells:
                            if first_row:
                                html_table += "<th>" + html.escape(cell.text) + "</th>"
                            else:
                                html_table += "<td>" + html.escape(cell.text) + "</td>"
                        html_table += "</tr>"
                        first_row = False
                    html_table += "</table></body></html>"
                    md_content += "\n" + self._convert(html_table).text_content.strip() + "\n"

                # Text areas
                elif shape.has_text_frame:
                    if shape == title:
                        md_content += "# " + shape.text.lstrip() + " "
                    else:
                        md_content += shape.text + " "

            md_content = md_content.strip()

            if slide.has_notes_slide:
                md_content += "\n\n### Notes:\n"
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    md_content += notes_frame.text
                md_content = md_content.strip()

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )

    def _is_picture(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def _is_table(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            return True
        return False


class FileConversionException(Exception):
    pass

class UnsupportedFormatException(Exception):
    pass


class MarkdownConverter:
    """(In preview) An extremely simple text-based document reader, suitable for LLM use.
    This reader will convert common file-types or webpages to Markdown."""
    
    usage_statistics = None
    
    def __init__(
        self,
        usage_statistics: Optional[UsageStatistics] = None,
        requests_session: Optional[requests.Session] = None,
    ):
        if requests_session is None:
            self._requests_session = requests.Session()
        else:
            self._requests_session = requests_session
        if usage_statistics is not None:
            self.usage_statistics = usage_statistics

        self._page_converters: List[DocumentConverter] = []

        # Register converters for successful browsing operations
        # Later registrations are tried first / take higher priority than earlier registrations
        # To this end, the most specific converters should appear below the most generic converters
        self.register_page_converter(XmlConverter())
        self.register_page_converter(YouTubeConverter())
        self.register_page_converter(DocxConverter())
        self.register_page_converter(XlsxConverter())
        self.register_page_converter(PptxConverter())
        # self.register_page_converter(ImageConverter())
        self.register_page_converter(PdfConverter())
        if self.usage_statistics is not None:
            self.register_page_converter(AudioConverter(self.usage_statistics))
        else:
            self.register_page_converter(AudioConverter())
        self.register_page_converter(HtmlConverter())
        self.register_page_converter(PlainTextConverter())

    def convert(self, source, **kwargs):
        """
        Args:
            - source: can be a string representing a path or url, or a requests.response object
            - extension: specifies the file extension to use when interpreting the file. If None, infer from source (path, uri, content-type, etc.)
        """
        # Local path or url
        if isinstance(source, str):
            if source.startswith("http://") or source.startswith("https://") or source.startswith("file://"):
                return self.convert_url(source, **kwargs)
            else:
                return self.convert_local(source, **kwargs)
        # Request response
        elif isinstance(source, requests.Response):
            return self.convert_response(source, **kwargs)

    def convert_local(self, path, **kwargs):
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Get extension alternatives from the path and puremagic
        base, ext = os.path.splitext(path)
        self._append_ext(extensions, ext)
        self._append_ext(extensions, self._guess_ext_magic(path))

        # Convert
        return self._convert(path, extensions, **kwargs)

    def convert_url(self, url, **kwargs):
        # Send a HTTP request to the URL
        user_agent = "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/119.0.0.0 Safari/537.36 Edg/119.0.0.0"
        response = self._requests_session.get(url, stream=True, headers={"User-Agent": user_agent})
        response.raise_for_status()
        return self.convert_response(response, **kwargs)

    def convert_response(self, response, **kwargs):
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []
        encoding = response.encoding
        # Guess from the mimetype
        content_type = response.headers.get("content-type", "").split(";")[0]
        self._append_ext(extensions, mimetypes.guess_extension(content_type))

        # Read the content disposition if there is one
        content_disposition = response.headers.get("content-disposition", "")
        m = re.search(r"filename=([^;]+)", content_disposition)
        if m:
            base, ext = os.path.splitext(m.group(1).strip("\"'"))
            self._append_ext(extensions, ext)

        # Read from the extension from the path
        base, ext = os.path.splitext(urlparse(response.url).path)
        self._append_ext(extensions, ext)

        # Save the file locally to a temporary file. It will be deleted before this method exits
        handle, temp_path = tempfile.mkstemp()
        fh = os.fdopen(handle, "wb")
        result = None
        try:
            # Download the file
            for chunk in response.iter_content(chunk_size=512):
                fh.write(chunk)
            fh.close()

            # Use puremagic to check for more extension options
            self._append_ext(extensions, self._guess_ext_magic(temp_path))

            # Convert
            result = self._convert(temp_path, extensions, url=response.url, encoding=encoding)

        # Clean up
        finally:
            try:
                fh.close()
            except:
                pass
            os.unlink(temp_path)

        return result

    def _convert(self, local_path, extensions, **kwargs):
        error_trace = ""
        for ext in extensions:
            for converter in self._page_converters:
                _kwargs = copy.deepcopy(kwargs)
                _kwargs.update({"file_extension": ext})
                # If we hit an error log it and keep trying
                try:
                    res = converter.convert(local_path, **_kwargs)
                    if res is not None:
                        # Normalize the content
                        res.text_content = "\n".join([line.rstrip() for line in re.split(r"\r?\n", res.text_content)])
                        res.text_content = re.sub(r"\n{3,}", "\n\n", res.text_content)

                        return res
                except Exception:
                    error_trace = ("\n\n" + traceback.format_exc()).strip()


        # If we got this far without success, report any exceptions
        if len(error_trace) > 0:
            raise FileConversionException(
                f"Could not convert '{local_path}' to Markdown. File type was recognized as {extensions}. While converting the file, the following error was encountered:\n\n{error_trace}"
            )

        # Nothing can handle it!
        # raise UnsupportedFormatException(
        #     f"Could not convert '{local_path}' to Markdown. The formats {extensions} are not supported."
        # )
        res = PlainTextConverter().convert(local_path, **kwargs)
        return res

    def _append_ext(self, extensions, ext):
        """Append a unique non-None, non-empty extension to a list of extensions."""
        if ext is None:
            return
        ext = ext.strip()
        if ext == "":
            return
        # if ext not in extensions:
        if True:
            extensions.append(ext)

    def _guess_ext_magic(self, path):
        """Use puremagic (a Python implementation of libmagic) to guess a file's extension based on the first few bytes."""
        # Use puremagic to guess
        try:
            guesses = puremagic.magic_file(path)
            if len(guesses) > 0:
                ext = guesses[0].extension.strip()
                if len(ext) > 0:
                    return ext
        except FileNotFoundError:
            pass
        except IsADirectoryError:
            pass
        except PermissionError:
            pass
        return None

    def register_page_converter(self, converter: DocumentConverter) -> None:
        """Register a page text converter."""
        self._page_converters.append(converter)

